# =============================================================================
# SCRIPT 07 - CREATE POWERPOINT PRESENTATION
# =============================================================================
# PURPOSE:
#   Generates a professional Master's-level PowerPoint presentation
#   based on the VaR project results. Uses actual data from the project
#   outputs (CSV files and chart images).
#
# REQUIRES: pip install python-pptx
# =============================================================================

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import csv
import math

# --- Paths ---
BASE_DIR    = os.path.dirname(os.path.dirname(__file__))
DATA_FOLDER = os.path.join(BASE_DIR, "DATA")
OUT_FOLDER  = os.path.join(BASE_DIR, "PRESENTATION")
OUT_FILE    = os.path.join(OUT_FOLDER, "VaR_Presentation_Master1_Finance.pptx")

CHART1 = os.path.join(DATA_FOLDER, "chart1_stock_performance.png")
CHART2 = os.path.join(DATA_FOLDER, "chart2_portfolio_value.png")
CHART3 = os.path.join(DATA_FOLDER, "chart3_return_distribution.png")
CHART4 = os.path.join(DATA_FOLDER, "chart4_var_comparison.png")

# =============================================================================
# LOAD PROJECT DATA (from CSVs — no invented numbers)
# =============================================================================

# --- VaR results ---
var_data = {}
with open(os.path.join(DATA_FOLDER, "var_results.csv")) as f:
    for row in csv.DictReader(f):
        var_data[row["Method"]] = float(row["VaR (%)"])

hist_var  = var_data["Historical VaR"]    # 1.7456
param_var = var_data["Parametric VaR"]   # 2.0295
mc_var    = var_data["Monte Carlo VaR"]  # 2.0422

# --- Portfolio value ---
with open(os.path.join(DATA_FOLDER, "portfolio_value.csv")) as f:
    pv_rows = list(csv.DictReader(f))

pv_vals  = [float(r["Portfolio_Value"]) for r in pv_rows]
pv_dates = [r["Date"] for r in pv_rows]
start_date  = pv_dates[0]          # 2020-01-03
end_date    = pv_dates[-1]         # 2024-12-30
final_val   = pv_vals[-1]          # ~18000
min_val     = min(pv_vals)
min_date    = pv_dates[pv_vals.index(min_val)]
max_val     = max(pv_vals)
total_ret   = (final_val - 10000) / 10000 * 100

# --- Portfolio returns ---
with open(os.path.join(DATA_FOLDER, "portfolio_returns.csv")) as f:
    pr_rows = list(csv.DictReader(f))

rets  = [float(r["Portfolio_Return"]) for r in pr_rows]
n_obs = len(rets)
mean_r = sum(rets) / n_obs
std_r  = math.sqrt(sum((r - mean_r) ** 2 for r in rets) / (n_obs - 1))
best_day  = max(rets) * 100
worst_day = min(rets) * 100
ann_ret   = ((1 + mean_r) ** 252 - 1) * 100
ann_vol   = std_r * math.sqrt(252) * 100
sharpe    = (mean_r / std_r) * math.sqrt(252)

print("=" * 60)
print("PROJECT DATA LOADED")
print(f"  Period         : {start_date} → {end_date}")
print(f"  Observations   : {n_obs} trading days")
print(f"  Final value    : ${final_val:,.2f}")
print(f"  Total return   : {total_ret:.2f}%")
print(f"  Worst day      : {worst_day:.2f}%  |  Best day: {best_day:.2f}%")
print(f"  Ann. return    : {ann_ret:.2f}%  |  Ann. vol: {ann_vol:.2f}%")
print(f"  Sharpe ratio   : {sharpe:.2f}")
print(f"  Historical VaR : {hist_var:.4f}%")
print(f"  Parametric VaR : {param_var:.4f}%")
print(f"  Monte Carlo VaR: {mc_var:.4f}%")
print("=" * 60)

# =============================================================================
# COLOUR PALETTE & HELPER FUNCTIONS
# =============================================================================

# Professional finance palette
NAVY        = RGBColor(0x0D, 0x2B, 0x55)   # Deep navy — title background
DARK_BLUE   = RGBColor(0x1A, 0x48, 0x8A)   # Section headers
ACCENT_BLUE = RGBColor(0x21, 0x96, 0xF3)   # Highlights / bullets
GOLD        = RGBColor(0xC8, 0xA0, 0x32)   # Accent / emphasis
RED         = RGBColor(0xC0, 0x39, 0x2B)   # Risk / VaR warnings
GREEN       = RGBColor(0x1A, 0x7A, 0x3C)   # Positive returns
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF4, 0xF6, 0xF9)
MID_GREY    = RGBColor(0x7F, 0x8C, 0x8D)
BLACK       = RGBColor(0x0A, 0x0A, 0x0A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(14), font_color=BLACK, bold=False,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size     = font_size
    run.font.color.rgb = font_color
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.name     = font_name
    return txBox


def add_para(tf, text, font_size=Pt(13), font_color=BLACK, bold=False,
             italic=False, align=PP_ALIGN.LEFT, space_before=Pt(4), font_name="Calibri"):
    """Add a paragraph to an existing text frame."""
    from pptx.util import Pt as pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size      = font_size
    run.font.color.rgb = font_color
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.name      = font_name
    return p


def header_bar(slide, title_text, subtitle_text=""):
    """Add a standard navy header bar to a content slide."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=NAVY)
    add_text(slide, title_text, 0.35, 0.1, 12.5, 0.65,
             font_size=Pt(26), font_color=WHITE, bold=True,
             align=PP_ALIGN.LEFT, font_name="Calibri Light")
    if subtitle_text:
        add_text(slide, subtitle_text, 0.35, 0.72, 12.5, 0.35,
                 font_size=Pt(13), font_color=ACCENT_BLUE,
                 align=PP_ALIGN.LEFT, font_name="Calibri")


def footer(slide, text="Value at Risk Analysis — Master 1 Finance | IAE Metz — 2025"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=NAVY)
    add_text(slide, text, 0.2, 7.21, 13.0, 0.28,
             font_size=Pt(9), font_color=RGBColor(0xAA, 0xBB, 0xCC),
             align=PP_ALIGN.CENTER)


def kpi_box(slide, label, value, l, t, w=2.8, h=1.3,
            bg=DARK_BLUE, val_color=WHITE, lbl_color=ACCENT_BLUE):
    """Render a KPI card (label + big value)."""
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_text(slide, label, l + 0.1, t + 0.08, w - 0.2, 0.4,
             font_size=Pt(11), font_color=lbl_color, bold=False,
             align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(slide, value, l + 0.1, t + 0.45, w - 0.2, 0.65,
             font_size=Pt(24), font_color=val_color, bold=True,
             align=PP_ALIGN.CENTER, font_name="Calibri Light")


# =============================================================================
# SLIDE 1 — TITLE SLIDE
# =============================================================================

prs = new_prs()

sl1 = blank_slide(prs)
fill_bg(sl1, NAVY)

# Gold accent bar top
add_rect(sl1, 0, 0, 13.33, 0.12, fill_color=GOLD)
# Gold accent bar bottom
add_rect(sl1, 0, 7.38, 13.33, 0.12, fill_color=GOLD)

# Decorative side panel
add_rect(sl1, 0, 0.12, 3.2, 7.26, fill_color=DARK_BLUE)
add_rect(sl1, 3.2, 0.12, 0.05, 7.26, fill_color=GOLD)

# Left panel labels
add_text(sl1, "MASTER 1\nFINANCE", 0.2, 0.5, 2.8, 1.2,
         font_size=Pt(13), font_color=GOLD, bold=True,
         align=PP_ALIGN.CENTER, font_name="Calibri Light")
add_text(sl1, "IAE METZ\n2024–2025", 0.2, 1.6, 2.8, 0.9,
         font_size=Pt(11), font_color=RGBColor(0xAA, 0xBB, 0xCC),
         align=PP_ALIGN.CENTER, font_name="Calibri Light")

# Divider line on left panel
add_rect(sl1, 0.4, 2.4, 2.4, 0.03, fill_color=GOLD)

# Left panel — VaR summary snapshot
add_text(sl1, "VaR Results (95%, 1-day)", 0.15, 2.6, 2.9, 0.4,
         font_size=Pt(10), font_color=ACCENT_BLUE, bold=True,
         align=PP_ALIGN.CENTER)
for i, (label, val) in enumerate([
    ("Historical", f"{hist_var:.4f}%"),
    ("Parametric", f"{param_var:.4f}%"),
    ("Monte Carlo", f"{mc_var:.4f}%"),
]):
    y = 3.05 + i * 0.62
    add_rect(sl1, 0.25, y, 2.7, 0.52, fill_color=RGBColor(0x0A, 0x1E, 0x40))
    add_text(sl1, label, 0.3, y + 0.02, 1.4, 0.28,
             font_size=Pt(9), font_color=MID_GREY, align=PP_ALIGN.LEFT)
    add_text(sl1, val, 0.3, y + 0.22, 2.5, 0.28,
             font_size=Pt(16), font_color=WHITE, bold=True, align=PP_ALIGN.LEFT)

add_text(sl1, f"Portfolio: ${final_val:,.0f} (+{total_ret:.0f}%)", 0.2, 5.0, 2.8, 0.4,
         font_size=Pt(10), font_color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(sl1, f"{start_date[:4]}–{end_date[:4]} | {n_obs} trading days", 0.2, 5.38, 2.8, 0.35,
         font_size=Pt(9), font_color=MID_GREY, align=PP_ALIGN.CENTER)

# Main title
add_text(sl1, "Value at Risk Analysis", 3.5, 1.2, 9.5, 1.2,
         font_size=Pt(42), font_color=WHITE, bold=True,
         align=PP_ALIGN.LEFT, font_name="Calibri Light")
add_text(sl1, "of a US Equity Portfolio", 3.5, 2.25, 9.5, 0.9,
         font_size=Pt(36), font_color=ACCENT_BLUE, bold=False,
         align=PP_ALIGN.LEFT, font_name="Calibri Light")

add_rect(sl1, 3.5, 3.1, 8.5, 0.04, fill_color=GOLD)

add_text(sl1, "Three-Method Approach  ·  Historical  ·  Parametric  ·  Monte Carlo",
         3.5, 3.25, 9.5, 0.5,
         font_size=Pt(14), font_color=GOLD, bold=False,
         align=PP_ALIGN.LEFT, font_name="Calibri Light")

add_text(sl1,
         "LVMH (Luxury)  ·  Johnson & Johnson (Healthcare)  ·  Lockheed Martin (Defense)\n"
         "NextEra Energy (Renewables)  ·  Apple (Technology)",
         3.5, 3.85, 9.5, 0.8,
         font_size=Pt(12), font_color=RGBColor(0xAA, 0xBB, 0xCC),
         align=PP_ALIGN.LEFT, font_name="Calibri Light")

add_text(sl1, "January 2020 – December 2024", 3.5, 4.75, 5.0, 0.45,
         font_size=Pt(13), font_color=WHITE, bold=False,
         align=PP_ALIGN.LEFT)

add_text(sl1, "Equal-Weighted Portfolio  ·  Initial Capital: $10,000",
         3.5, 5.15, 7.0, 0.4,
         font_size=Pt(12), font_color=MID_GREY,
         align=PP_ALIGN.LEFT)

add_text(sl1, "Presented by:", 3.5, 6.1, 3.0, 0.3,
         font_size=Pt(10), font_color=MID_GREY, align=PP_ALIGN.LEFT)
add_text(sl1, "Master 1 Finance — IAE Metz", 3.5, 6.4, 7.0, 0.4,
         font_size=Pt(12), font_color=WHITE, bold=False, align=PP_ALIGN.LEFT)

# =============================================================================
# SLIDE 2 — CONTEXT & OBJECTIVES
# =============================================================================

sl2 = blank_slide(prs)
fill_bg(sl2, LIGHT_GREY)
header_bar(sl2, "Context & Objectives",
           "Why measure risk?  What does this project answer?")
footer(sl2)

# Left column — context
add_rect(sl2, 0.3, 1.25, 6.0, 5.65, fill_color=WHITE)
add_rect(sl2, 0.3, 1.25, 0.08, 5.65, fill_color=DARK_BLUE)  # accent bar

add_text(sl2, "Financial Context", 0.55, 1.35, 5.5, 0.45,
         font_size=Pt(16), font_color=DARK_BLUE, bold=True)

ctx_tf = sl2.shapes.add_textbox(Inches(0.55), Inches(1.9), Inches(5.6), Inches(4.8)).text_frame
ctx_tf.word_wrap = True
p0 = ctx_tf.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT

bullets_ctx = [
    ("CONTEXT", True, DARK_BLUE, Pt(12)),
    ("Since the 2008 financial crisis, VaR has become the standard regulatory "
     "risk metric under Basel III/IV. Every major bank and asset manager is "
     "required to report daily VaR figures.", False, BLACK, Pt(12)),
    ("", False, BLACK, Pt(6)),
    ("PORTFOLIO RATIONALE", True, DARK_BLUE, Pt(12)),
    ("We build a 5-stock US equity portfolio across uncorrelated sectors: "
     "luxury goods, healthcare, defense, renewable energy, and technology. "
     "The goal is to observe how sectoral diversification affects risk.", False, BLACK, Pt(12)),
    ("", False, BLACK, Pt(6)),
    ("PERIOD CHOICE", True, DARK_BLUE, Pt(12)),
    ("January 2020 – December 2024 deliberately covers extreme market events: "
     "the COVID crash (March 2020), the post-COVID rally, and the 2022 rate-hike "
     "bear market. This stress-tests our VaR models against real shocks.", False, BLACK, Pt(12)),
]

first = True
for (txt, bold, color, size) in bullets_ctx:
    if first:
        run = p0.add_run()
        run.text = txt
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size
        run.font.name = "Calibri"
        first = False
    else:
        p = ctx_tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = txt
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size
        run.font.name = "Calibri"

# Right column — objectives
add_rect(sl2, 6.7, 1.25, 6.3, 5.65, fill_color=WHITE)
add_rect(sl2, 6.7, 1.25, 0.08, 5.65, fill_color=GOLD)

add_text(sl2, "Research Objectives", 6.95, 1.35, 5.8, 0.45,
         font_size=Pt(16), font_color=DARK_BLUE, bold=True)

objectives = [
    ("01", "Construct a diversified US equity portfolio across 5 distinct sectors using 5 years of daily market data."),
    ("02", "Compute 1-day VaR at the 95% confidence level using three independent methods."),
    ("03", "Compare and critically interpret the three VaR estimates — do they converge or diverge, and why?"),
    ("04", "Assess what the VaR figures mean in dollar terms for a $10,000 invested portfolio."),
]

for i, (num, txt) in enumerate(objectives):
    y = 1.95 + i * 1.2
    add_rect(sl2, 6.95, y, 0.55, 0.55, fill_color=DARK_BLUE)
    add_text(sl2, num, 6.97, y + 0.05, 0.5, 0.45,
             font_size=Pt(18), font_color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl2, txt, 7.6, y, 5.15, 0.9,
             font_size=Pt(12), font_color=BLACK, wrap=True)

# =============================================================================
# SLIDE 3 — PORTFOLIO CONSTRUCTION
# =============================================================================

sl3 = blank_slide(prs)
fill_bg(sl3, LIGHT_GREY)
header_bar(sl3, "Portfolio Construction",
           "5 stocks · 5 sectors · Equal weights · $10,000 initial capital")
footer(sl3)

# Portfolio composition table
stocks = [
    ("LVMUY", "LVMH Moët Hennessy Louis Vuitton", "Luxury Goods", "Europe / Global", "20%"),
    ("JNJ",   "Johnson & Johnson",                  "Healthcare",   "USA",             "20%"),
    ("LMT",   "Lockheed Martin",                    "Defense",      "USA",             "20%"),
    ("NEE",   "NextEra Energy",                     "Renewable Energy","USA",           "20%"),
    ("AAPL",  "Apple Inc.",                          "Technology",   "USA",             "20%"),
]

col_headers = ["Ticker", "Company", "Sector", "Geography", "Weight"]
col_widths  = [1.0, 3.5, 2.2, 1.7, 0.9]
col_x       = [0.3]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# Header row
add_rect(sl3, 0.3, 1.25, sum(col_widths), 0.42, fill_color=NAVY)
for j, (hdr, x, w) in enumerate(zip(col_headers, col_x, col_widths)):
    add_text(sl3, hdr, x + 0.05, 1.28, w - 0.1, 0.36,
             font_size=Pt(12), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

row_colors = [WHITE, LIGHT_GREY]
for i, (tkr, name, sector, geo, wgt) in enumerate(stocks):
    y = 1.67 + i * 0.52
    bg = row_colors[i % 2]
    add_rect(sl3, 0.3, y, sum(col_widths), 0.52, fill_color=bg)
    vals = [tkr, name, sector, geo, wgt]
    for j, (val, x, w) in enumerate(zip(vals, col_x, col_widths)):
        bold = (j == 0)
        col  = DARK_BLUE if j == 0 else BLACK
        add_text(sl3, val, x + 0.05, y + 0.08, w - 0.1, 0.36,
                 font_size=Pt(12), font_color=col, bold=bold,
                 align=PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT)

# Rationale box below table
add_rect(sl3, 0.3, 4.3, 12.7, 2.65, fill_color=WHITE)
add_rect(sl3, 0.3, 4.3, 0.08, 2.65, fill_color=ACCENT_BLUE)

add_text(sl3, "Construction Rationale", 0.55, 4.38, 12.0, 0.4,
         font_size=Pt(15), font_color=DARK_BLUE, bold=True)

rationale = [
    "DIVERSIFICATION LOGIC  — The five sectors have low return correlations. "
    "Defence (LMT) and healthcare (JNJ) tend to hold value during equity sell-offs, "
    "while technology (AAPL) and luxury (LVMUY) drive growth in risk-on environments. "
    "Renewables (NEE) provide exposure to long-term structural growth with low sensitivity to business cycles.",

    "EQUAL WEIGHTING  — Each stock receives exactly 20% of capital. This avoids concentration bias and makes "
    "the portfolio transparent and reproducible. In academic finance, equal weighting is a well-established "
    "benchmark that often outperforms optimised portfolios out-of-sample.",
]

rat_box = sl3.shapes.add_textbox(Inches(0.55), Inches(4.78), Inches(12.5), Inches(2.0))
rat_tf  = rat_box.text_frame
rat_tf.word_wrap = True
first2 = True
for txt in rationale:
    if first2:
        p = rat_tf.paragraphs[0]
        first2 = False
    else:
        p = rat_tf.add_paragraph()
        p.space_before = Pt(6)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(11.5)
    run.font.color.rgb = BLACK
    run.font.name = "Calibri"

# =============================================================================
# SLIDE 4 — STOCK PERFORMANCE (Chart 1)
# =============================================================================

sl4 = blank_slide(prs)
fill_bg(sl4, LIGHT_GREY)
header_bar(sl4, "Individual Stock Performance",
           "Normalized to base 100 — January 2020 to December 2024")
footer(sl4)

# Chart image (left 2/3)
sl4.shapes.add_picture(CHART1, Inches(0.2), Inches(1.15), Inches(8.8), Inches(5.1))

# Interpretation panel (right 1/3)
add_rect(sl4, 9.2, 1.15, 3.93, 5.1, fill_color=WHITE)
add_rect(sl4, 9.2, 1.15, 0.07, 5.1, fill_color=DARK_BLUE)

add_text(sl4, "Key Observations", 9.4, 1.25, 3.65, 0.4,
         font_size=Pt(14), font_color=DARK_BLUE, bold=True)

obs_items = [
    ("AAPL — Dominant performer", DARK_BLUE, True),
    ("Apple delivered the strongest cumulative gain, driven by product cycle expansion and multiple re-rating.", BLACK, False),
    ("", BLACK, False),
    ("LMT — Defensive outperformer", DARK_BLUE, True),
    ("Lockheed Martin proved resilient through the 2022 bear market as defence budgets surged post-Ukraine.", BLACK, False),
    ("", BLACK, False),
    ("NEE — High beta to rates", DARK_BLUE, True),
    ("NextEra lost ground significantly in 2022–23 as rising rates increased its cost of capital (utility sector sensitivity).", BLACK, False),
    ("", BLACK, False),
    ("LVMUY — Volatile recovery", DARK_BLUE, True),
    ("LVMH recovered from COVID at pace but then corrected in 2024 as China consumption disappointed.", BLACK, False),
    ("", BLACK, False),
    ("JNJ — Stable anchor", DARK_BLUE, True),
    ("Johnson & Johnson showed the lowest volatility, consistent with its defensive healthcare profile.", BLACK, False),
]

obs_box = sl4.shapes.add_textbox(Inches(9.4), Inches(1.75), Inches(3.65), Inches(4.3))
obs_tf  = obs_box.text_frame
obs_tf.word_wrap = True
first3  = True
for (txt, color, bold) in obs_items:
    if first3:
        p = obs_tf.paragraphs[0]; first3 = False
    else:
        p = obs_tf.add_paragraph()
        p.space_before = Pt(2)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(10.5)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"

# =============================================================================
# SLIDE 5 — PORTFOLIO VALUE EVOLUTION (Chart 2)
# =============================================================================

sl5 = blank_slide(prs)
fill_bg(sl5, LIGHT_GREY)
header_bar(sl5, "Portfolio Value Evolution",
           f"$10,000 invested equally across 5 stocks · Final value: ${final_val:,.0f}")
footer(sl5)

# Chart image
sl5.shapes.add_picture(CHART2, Inches(0.2), Inches(1.15), Inches(8.8), Inches(4.3))

# KPI cards below chart
kpis = [
    ("Initial Capital", "$10,000"),
    ("Final Value", f"${final_val:,.0f}"),
    ("Total Return", f"+{total_ret:.1f}%"),
    ("Ann. Return", f"+{ann_ret:.1f}%"),
    ("Worst Day", f"{worst_day:.2f}%"),
    ("Ann. Volatility", f"{ann_vol:.1f}%"),
]
kpi_colors = [DARK_BLUE, GREEN, GREEN, GREEN, RED, DARK_BLUE]
for i, ((lbl, val), bg) in enumerate(zip(kpis, kpi_colors)):
    kpi_box(sl5, lbl, val, 0.2 + i * 1.49, 5.57, w=1.42, h=1.18, bg=bg)

# Interpretation panel
add_rect(sl5, 9.2, 1.15, 3.93, 4.3, fill_color=WHITE)
add_rect(sl5, 9.2, 1.15, 0.07, 4.3, fill_color=GOLD)

add_text(sl5, "Portfolio Narrative", 9.4, 1.25, 3.65, 0.4,
         font_size=Pt(14), font_color=DARK_BLUE, bold=True)

narr_items = [
    ("PHASE 1 — COVID Crash (Feb–Mar 2020)", DARK_BLUE, True),
    (f"Portfolio dropped sharply to ${min_val:,.0f} (on {min_date}). The equal-weight structure meant all sectors were hit, but JNJ and LMT cushioned the fall.", BLACK, False),
    ("", BLACK, False),
    ("PHASE 2 — Bull Market (2020–2021)", DARK_BLUE, True),
    ("Strong recovery driven by AAPL and LVMUY. The portfolio exceeded the initial $10,000 by mid-2020 and continued to grow through 2021.", BLACK, False),
    ("", BLACK, False),
    ("PHASE 3 — Rate Hike Bear (2022)", DARK_BLUE, True),
    ("The Fed's aggressive tightening compressed growth valuations. AAPL and NEE dragged the portfolio down. LMT bucked the trend.", BLACK, False),
    ("", BLACK, False),
    ("PHASE 4 — Recovery (2023–2024)", DARK_BLUE, True),
    (f"AI-driven rally in tech and stabilising inflation pushed the portfolio to its final value of ${final_val:,.0f}, a gain of {total_ret:.1f}% over the full period.", BLACK, False),
]

narr_box = sl5.shapes.add_textbox(Inches(9.4), Inches(1.72), Inches(3.65), Inches(3.65))
narr_tf  = narr_box.text_frame
narr_tf.word_wrap = True
first4   = True
for (txt, color, bold) in narr_items:
    if first4:
        p = narr_tf.paragraphs[0]; first4 = False
    else:
        p = narr_tf.add_paragraph()
        p.space_before = Pt(2)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(10)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"

# =============================================================================
# SLIDE 6 — METHODOLOGY: WHAT IS VaR & HOW WE COMPUTED IT
# =============================================================================

sl6 = blank_slide(prs)
fill_bg(sl6, LIGHT_GREY)
header_bar(sl6, "Value at Risk — Methodology",
           "Three independent estimation methods · 95% confidence · 1-day horizon")
footer(sl6)

# Central definition box
add_rect(sl6, 0.3, 1.25, 12.7, 1.15, fill_color=DARK_BLUE)
add_text(sl6,
         "Definition — VaR (95%, 1-day) answers: "
         "\"What is the maximum loss this portfolio will suffer on a typical bad day, "
         "19 times out of 20?\"",
         0.45, 1.3, 12.4, 1.05,
         font_size=Pt(15), font_color=WHITE, bold=False,
         align=PP_ALIGN.CENTER, font_name="Calibri Light")

# Three method cards
method_data = [
    ("01", "Historical VaR",
     "Uses the actual distribution of past returns.",
     [
         f"· Take all {n_obs} daily returns (Jan 2020 – Dec 2024)",
         "· Sort returns from worst to best",
         "· Read off the 5th percentile directly",
         "→ No distributional assumption",
         "→ Captures real market events (COVID, 2022 crash)",
         "→ Limitation: future may differ from past",
     ],
     f"{hist_var:.4f}%"),
    ("02", "Parametric VaR",
     "Assumes portfolio returns follow a Normal distribution.",
     [
         f"· Mean daily return: {mean_r*100:.4f}%",
         f"· Daily std deviation: {std_r*100:.4f}%",
         "· Z-score at 5% tail: −1.645",
         "→ Formula: VaR = −(μ + z·σ)",
         "→ Fast and analytically tractable",
         "→ Risk: underestimates fat tails",
     ],
     f"{param_var:.4f}%"),
    ("03", "Monte Carlo VaR",
     "Simulates 10,000 possible future return scenarios.",
     [
         "· Draw 10,000 random returns from N(μ, σ²)",
         "· Same parameters as parametric method",
         "· Take the 5th percentile of simulated outcomes",
         "→ Probabilistic view of the loss distribution",
         "→ Very close to parametric (same inputs)",
         "→ More flexible for complex portfolios",
     ],
     f"{mc_var:.4f}%"),
]

card_x = [0.3, 4.55, 8.8]
for i, (num, title, subtitle, bullets, result) in enumerate(method_data):
    x = card_x[i]
    w = 4.1
    # Card background
    add_rect(sl6, x, 2.55, w, 4.2, fill_color=WHITE)
    # Top colour strip
    add_rect(sl6, x, 2.55, w, 0.08, fill_color=[DARK_BLUE, GOLD, ACCENT_BLUE][i])
    # Number badge
    add_rect(sl6, x + 0.12, 2.65, 0.5, 0.5, fill_color=NAVY)
    add_text(sl6, num, x + 0.13, 2.66, 0.47, 0.48,
             font_size=Pt(16), font_color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    # Title
    add_text(sl6, title, x + 0.72, 2.67, w - 0.85, 0.42,
             font_size=Pt(14), font_color=DARK_BLUE, bold=True)
    add_text(sl6, subtitle, x + 0.15, 3.12, w - 0.25, 0.38,
             font_size=Pt(10.5), font_color=MID_GREY, italic=True)

    # Bullets
    blt_box = sl6.shapes.add_textbox(Inches(x + 0.15), Inches(3.52), Inches(w - 0.25), Inches(2.45))
    blt_tf  = blt_box.text_frame
    blt_tf.word_wrap = True
    first5  = True
    for b in bullets:
        bold = b.startswith("→")
        color = DARK_BLUE if bold else BLACK
        if first5:
            p = blt_tf.paragraphs[0]; first5 = False
        else:
            p = blt_tf.add_paragraph()
            p.space_before = Pt(2)
        run = p.add_run()
        run.text = b
        run.font.size = Pt(10.5)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Calibri"

    # Result badge
    add_rect(sl6, x + 0.1, 6.35, w - 0.2, 0.32, fill_color=NAVY)
    add_text(sl6, f"Result: {result}", x + 0.1, 6.36, w - 0.2, 0.3,
             font_size=Pt(12), font_color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 7 — RETURN DISTRIBUTION (Chart 3)
# =============================================================================

sl7 = blank_slide(prs)
fill_bg(sl7, LIGHT_GREY)
header_bar(sl7, "Daily Return Distribution",
           "Histogram of actual returns vs. theoretical normal curve · Historical VaR threshold")
footer(sl7)

sl7.shapes.add_picture(CHART3, Inches(0.2), Inches(1.15), Inches(8.8), Inches(5.1))

add_rect(sl7, 9.2, 1.15, 3.93, 5.1, fill_color=WHITE)
add_rect(sl7, 9.2, 1.15, 0.07, 5.1, fill_color=RED)

add_text(sl7, "Reading the Chart", 9.4, 1.25, 3.65, 0.4,
         font_size=Pt(14), font_color=DARK_BLUE, bold=True)

dist_items = [
    ("Central mass — low risk days", DARK_BLUE, True),
    ("The bulk of daily returns cluster tightly around 0%, confirming that most trading days are unremarkable for a diversified portfolio.", BLACK, False),
    ("", BLACK, False),
    ("Left tail — the loss zone", RED, True),
    (f"The red vertical line marks the Historical VaR threshold at −{hist_var:.2f}%. "
     f"On 5% of trading days (≈ 63 days over 5 years), losses exceeded this level.", BLACK, False),
    ("", BLACK, False),
    ("Normal curve overlay", DARK_BLUE, True),
    (f"The orange curve is the theoretical normal distribution fitted to μ = {mean_r*100:.3f}% "
     f"and σ = {std_r*100:.3f}%. The actual histogram is slightly leptokurtic "
     "(more peaked, fatter tails) — typical of equity returns.", BLACK, False),
    ("", BLACK, False),
    ("Key statistic", DARK_BLUE, True),
    (f"Worst single day: {worst_day:.2f}%  |  Best: +{best_day:.2f}%", RED, True),
]

dist_box = sl7.shapes.add_textbox(Inches(9.4), Inches(1.72), Inches(3.65), Inches(4.4))
dist_tf  = dist_box.text_frame
dist_tf.word_wrap = True
first6   = True
for (txt, color, bold) in dist_items:
    if first6:
        p = dist_tf.paragraphs[0]; first6 = False
    else:
        p = dist_tf.add_paragraph()
        p.space_before = Pt(2)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(10.5)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"

# =============================================================================
# SLIDE 8 — VaR RESULTS & INTERPRETATION
# =============================================================================

sl8 = blank_slide(prs)
fill_bg(sl8, LIGHT_GREY)
header_bar(sl8, "VaR Results — 95% Confidence, 1-Day Horizon",
           "Comparing three methods on the same portfolio")
footer(sl8)

# Three large VaR cards
card_data_8 = [
    ("Historical VaR", f"{hist_var:.4f}%", f"−${hist_var/100*final_val:,.0f}",
     "#E53935", "Most conservative estimate.",
     "Reads directly from the actual distribution of past returns. "
     f"On the 5th worst percentile of the {n_obs} observed trading days, "
     f"the portfolio lost at least {hist_var:.2f}% — or about ${hist_var/100*final_val:,.0f} "
     f"on the current portfolio value of ${final_val:,.0f}.",
     "Lower than parametric → the actual return distribution has a less extreme left tail than a normal curve predicts. Diversification worked."),
    ("Parametric VaR", f"{param_var:.4f}%", f"−${param_var/100*final_val:,.0f}",
     "#FB8C00", "Assumes Gaussian returns.",
     f"Derived analytically: VaR = −(μ + z·σ) with μ={mean_r*100:.4f}%, "
     f"σ={std_r*100:.4f}%, z=−1.645. "
     f"Gives {param_var:.2f}%, or about ${param_var/100*final_val:,.0f} in dollar terms.",
     "Higher than Historical → the normal distribution overestimates tail risk here. Parametric VaR is the more conservative parametric estimate."),
    ("Monte Carlo VaR", f"{mc_var:.4f}%", f"−${mc_var/100*final_val:,.0f}",
     "#8E24AA", "10,000 simulated scenarios.",
     f"Simulates 10,000 random draws from N({mean_r*100:.4f}%, {std_r*100:.4f}%²). "
     f"The 5th percentile of simulated outcomes gives {mc_var:.2f}%, "
     f"or ${mc_var/100*final_val:,.0f}.",
     f"Near-identical to Parametric ({mc_var:.4f}% vs {param_var:.4f}%). "
     "Expected: both use the same μ and σ. The 0.01% difference is pure sampling noise from 10,000 draws."),
]

cx = [0.3, 4.55, 8.8]
for i, (method, pct, dollar, color_hex, tagline, body, insight) in enumerate(card_data_8):
    x = cx[i]
    w = 4.1
    # Card
    add_rect(sl8, x, 1.25, w, 5.65, fill_color=WHITE)
    # Top bar
    c = RGBColor(int(color_hex[1:3],16), int(color_hex[3:5],16), int(color_hex[5:7],16))
    add_rect(sl8, x, 1.25, w, 0.12, fill_color=c)
    # Method name
    add_text(sl8, method, x + 0.15, 1.42, w - 0.25, 0.48,
             font_size=Pt(15), font_color=DARK_BLUE, bold=True)
    # Tagline
    add_text(sl8, tagline, x + 0.15, 1.87, w - 0.25, 0.32,
             font_size=Pt(10), font_color=MID_GREY, italic=True)
    # VaR percentage
    add_rect(sl8, x + 0.15, 2.24, w - 0.3, 0.88, fill_color=NAVY)
    add_text(sl8, pct, x + 0.15, 2.26, w - 0.3, 0.54,
             font_size=Pt(28), font_color=c, bold=True, align=PP_ALIGN.CENTER,
             font_name="Calibri Light")
    add_text(sl8, dollar + " on current portfolio",
             x + 0.15, 2.78, w - 0.3, 0.3,
             font_size=Pt(10), font_color=WHITE, align=PP_ALIGN.CENTER)
    # Body text
    body_box = sl8.shapes.add_textbox(Inches(x + 0.15), Inches(3.2), Inches(w - 0.25), Inches(1.55))
    body_tf  = body_box.text_frame
    body_tf.word_wrap = True
    p = body_tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(10.5)
    run.font.color.rgb = BLACK
    run.font.name = "Calibri"
    # Insight strip
    add_rect(sl8, x + 0.1, 4.8, w - 0.2, 2.0, fill_color=RGBColor(0xF0, 0xF4, 0xFA))
    add_rect(sl8, x + 0.1, 4.8, 0.06, 2.0, fill_color=c)
    ins_box = sl8.shapes.add_textbox(Inches(x + 0.25), Inches(4.85), Inches(w - 0.4), Inches(1.85))
    ins_tf  = ins_box.text_frame
    ins_tf.word_wrap = True
    p2 = ins_tf.paragraphs[0]
    run2 = p2.add_run()
    run2.text = insight
    run2.font.size = Pt(10)
    run2.font.color.rgb = DARK_BLUE
    run2.font.italic = True
    run2.font.name = "Calibri"

# =============================================================================
# SLIDE 9 — VaR COMPARISON CHART (Chart 4)
# =============================================================================

sl9 = blank_slide(prs)
fill_bg(sl9, LIGHT_GREY)
header_bar(sl9, "VaR Methods — Side-by-Side Comparison",
           "All three thresholds on the actual return distribution")
footer(sl9)

sl9.shapes.add_picture(CHART4, Inches(0.2), Inches(1.15), Inches(8.5), Inches(5.2))

# Summary table on right
add_rect(sl9, 8.9, 1.15, 4.23, 5.2, fill_color=WHITE)
add_rect(sl9, 8.9, 1.15, 0.07, 5.2, fill_color=DARK_BLUE)

add_text(sl9, "Comparison Summary", 9.1, 1.22, 3.9, 0.42,
         font_size=Pt(14), font_color=DARK_BLUE, bold=True)

# Comparison table
headers9 = ["Method", "VaR (%)", "$ Loss (on $18K)", "Assumption"]
rows9 = [
    ("Historical",   f"{hist_var:.4f}%",  f"−${hist_var/100*final_val:,.0f}",  "None (empirical)"),
    ("Parametric",   f"{param_var:.4f}%", f"−${param_var/100*final_val:,.0f}", "Normal dist."),
    ("Monte Carlo",  f"{mc_var:.4f}%",    f"−${mc_var/100*final_val:,.0f}",    "Normal dist."),
]

# Col x positions (relative within the right panel)
px = [8.92, 10.05, 11.25, 12.28]
pw = [1.1, 1.18, 0.97, 0.82]

add_rect(sl9, 8.9, 1.7, 4.23, 0.38, fill_color=NAVY)
for j, (h, x, w) in enumerate(zip(headers9, px, pw)):
    add_text(sl9, h, x, 1.72, w, 0.33,
             font_size=Pt(9), font_color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)

rc = [WHITE, LIGHT_GREY, WHITE]
bar_colors = ["#E53935", "#FB8C00", "#8E24AA"]
for i, (row, bg, bc) in enumerate(zip(rows9, rc, bar_colors)):
    y = 2.08 + i * 0.44
    rbg = WHITE if i % 2 == 0 else LIGHT_GREY
    add_rect(sl9, 8.9, y, 4.23, 0.44, fill_color=rbg)
    bc_rgb = RGBColor(int(bc[1:3],16), int(bc[3:5],16), int(bc[5:7],16))
    add_rect(sl9, 8.9, y, 0.07, 0.44, fill_color=bc_rgb)
    for j, (val, x, w) in enumerate(zip(row, px, pw)):
        bold = (j == 1)
        add_text(sl9, val, x, y + 0.06, w, 0.32,
                 font_size=Pt(9.5 if j > 0 else 10),
                 font_color=bc_rgb if j == 1 else BLACK,
                 bold=bold, align=PP_ALIGN.CENTER)

# Key takeaway
delta = param_var - hist_var
add_rect(sl9, 8.9, 3.4, 4.23, 0.06, fill_color=NAVY)

take_box = sl9.shapes.add_textbox(Inches(8.97), Inches(3.52), Inches(4.05), Inches(2.7))
take_tf  = take_box.text_frame
take_tf.word_wrap = True

takeaways = [
    ("KEY FINDING", DARK_BLUE, True),
    (f"Historical VaR ({hist_var:.2f}%) is {delta:.2f} percentage points lower "
     f"than the Parametric estimate ({param_var:.2f}%). This gap reveals that the portfolio's "
     "actual return distribution is less extreme in its left tail than a pure normal distribution predicts.", BLACK, False),
    ("", BLACK, False),
    ("WHY DOES THIS MATTER?", DARK_BLUE, True),
    ("A risk manager using only the parametric model would systematically over-provision capital "
     "against daily losses. The historical method, anchored to real data, provides a more accurate "
     "picture of this specific portfolio's behaviour.", BLACK, False),
    ("", BLACK, False),
    ("MONTE CARLO ≈ PARAMETRIC", DARK_BLUE, True),
    (f"The difference of {abs(mc_var - param_var):.4f}% between MC and Parametric is negligible — "
     "a direct consequence of identical input parameters. Monte Carlo's value lies in more complex, "
     "multi-asset scenarios.", BLACK, False),
]

first9 = True
for (txt, color, bold) in takeaways:
    if first9:
        p = take_tf.paragraphs[0]; first9 = False
    else:
        p = take_tf.add_paragraph()
        p.space_before = Pt(3)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(10)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"

# =============================================================================
# SLIDE 10 — KEY FINDINGS & FINANCIAL INTERPRETATION
# =============================================================================

sl10 = blank_slide(prs)
fill_bg(sl10, LIGHT_GREY)
header_bar(sl10, "Key Findings & Financial Interpretation",
           "What do the VaR numbers mean in practice?")
footer(sl10)

# Top row — three insight boxes
insights = [
    ("Diversification Benefit",
     f"The portfolio's Historical VaR of {hist_var:.2f}% is below the Parametric estimate, "
     f"confirming that cross-sector diversification reduces extreme tail losses beyond what "
     "a single-distribution model captures. Sectors moved counter-cyclically during 2020 and 2022.",
     DARK_BLUE, ACCENT_BLUE),
    ("Dollar Translation",
     f"On the worst 5% of trading days, the portfolio (valued at ${final_val:,.0f}) "
     f"is expected to lose at least ${hist_var/100*final_val:,.0f} (Historical) to "
     f"${param_var/100*final_val:,.0f} (Parametric). "
     "This translates the abstract percentage into a concrete risk budget.",
     GREEN, GREEN),
    ("Method Convergence",
     f"All three methods agree within a {mc_var - hist_var:.2f}% band. "
     "This convergence strengthens confidence in the risk estimate — no single method is an outlier. "
     "The Parametric and Monte Carlo methods converge because they share the same mean and volatility inputs.",
     GOLD, GOLD),
]

for i, (title, body, col, barcol) in enumerate(insights):
    x = 0.3 + i * 4.37
    add_rect(sl10, x, 1.25, 4.1, 2.85, fill_color=WHITE)
    add_rect(sl10, x, 1.25, 4.1, 0.08, fill_color=barcol)
    add_text(sl10, title, x + 0.15, 1.38, 3.8, 0.45,
             font_size=Pt(13), font_color=col, bold=True)
    body_box = sl10.shapes.add_textbox(Inches(x + 0.15), Inches(1.88), Inches(3.8), Inches(2.1))
    body_tf  = body_box.text_frame
    body_tf.word_wrap = True
    p = body_tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(11)
    run.font.color.rgb = BLACK
    run.font.name = "Calibri"

# Bottom — limitations
add_rect(sl10, 0.3, 4.2, 12.7, 2.65, fill_color=NAVY)
add_rect(sl10, 0.3, 4.2, 12.7, 0.06, fill_color=RED)

add_text(sl10, "Critical Limitations — What VaR Does NOT Tell You",
         0.5, 4.3, 12.3, 0.45,
         font_size=Pt(15), font_color=WHITE, bold=True)

limit_items = [
    ("VaR is NOT a worst-case loss. It is a threshold exceeded 5% of the time. "
     "On those 5% of days, losses can be far larger (this is measured by Expected Shortfall / CVaR)."),
    ("VaR assumes the past is representative of the future. "
     "A structural break (new pandemic, sovereign crisis) can invalidate historical distributions entirely."),
    ("The 1-day horizon ignores liquidity risk. "
     "Unwinding a $18K portfolio in a crisis may take days, during which losses can compound."),
    ("Equal weighting is a simplification. "
     "Optimal portfolio construction (mean-variance, risk parity) could reduce VaR further."),
]

lim_box = sl10.shapes.add_textbox(Inches(0.5), Inches(4.82), Inches(12.5), Inches(1.9))
lim_tf  = lim_box.text_frame
lim_tf.word_wrap = True
first10 = True
for txt in limit_items:
    if first10:
        p = lim_tf.paragraphs[0]; first10 = False
    else:
        p = lim_tf.add_paragraph()
        p.space_before = Pt(4)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "▸  " + txt
    run.font.size = Pt(10.5)
    run.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
    run.font.name = "Calibri"

# =============================================================================
# SLIDE 11 — CONCLUSION
# =============================================================================

sl11 = blank_slide(prs)
fill_bg(sl11, NAVY)

# Gold accent bars
add_rect(sl11, 0, 0, 13.33, 0.1, fill_color=GOLD)
add_rect(sl11, 0, 7.4, 13.33, 0.1, fill_color=GOLD)

# Right panel
add_rect(sl11, 9.8, 0.1, 3.53, 7.3, fill_color=DARK_BLUE)

# Title
add_text(sl11, "Conclusion", 0.5, 0.35, 9.0, 0.85,
         font_size=Pt(38), font_color=WHITE, bold=True,
         font_name="Calibri Light")
add_rect(sl11, 0.5, 1.15, 6.0, 0.05, fill_color=GOLD)

# Summary bullets
conclusions = [
    ("A well-diversified portfolio reduces VaR",
     f"By combining 5 low-correlated sectors, the portfolio's actual 1-day 95% VaR ({hist_var:.2f}%) "
     f"came in below what a normal-distribution model would predict ({param_var:.2f}%). "
     "This is a direct quantification of the diversification benefit."),
    ("All three methods converge — increasing confidence",
     f"Historical ({hist_var:.2f}%), Parametric ({param_var:.2f}%), and Monte Carlo ({mc_var:.2f}%) "
     f"align within a {mc_var - hist_var:.2f}% range. This convergence validates the robustness of the estimate "
     "and reduces model risk."),
    (f"Strong risk-adjusted performance over 5 years",
     f"The portfolio grew from $10,000 to ${final_val:,.0f} (+{total_ret:.1f}%) with an annualised "
     f"return of {ann_ret:.1f}% and annualised volatility of {ann_vol:.1f}%. "
     "VaR remained manageable throughout, including during the COVID shock."),
    ("VaR is a starting point, not the full risk picture",
     "Expected Shortfall (CVaR) should complement VaR to capture losses beyond the threshold. "
     "Dynamic rebalancing and correlation monitoring over time are essential for a live portfolio."),
]

for i, (hdr, body) in enumerate(conclusions):
    y = 1.3 + i * 1.4
    add_rect(sl11, 0.5, y, 0.06, 1.15, fill_color=GOLD)
    add_text(sl11, hdr, 0.7, y, 8.9, 0.42,
             font_size=Pt(13), font_color=GOLD, bold=True)
    bod_box = sl11.shapes.add_textbox(Inches(0.7), Inches(y + 0.42), Inches(8.9), Inches(0.85))
    bod_tf  = bod_box.text_frame
    bod_tf.word_wrap = True
    p = bod_tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(11.5)
    run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    run.font.name = "Calibri"

# Right panel content
add_text(sl11, "Final VaR Summary", 9.95, 0.5, 3.2, 0.42,
         font_size=Pt(14), font_color=GOLD, bold=True)
add_rect(sl11, 9.95, 0.92, 3.2, 0.04, fill_color=GOLD)

for i, (method, val, col_hex) in enumerate([
    ("Historical", f"{hist_var:.4f}%", "#E53935"),
    ("Parametric", f"{param_var:.4f}%", "#FB8C00"),
    ("Monte Carlo", f"{mc_var:.4f}%", "#8E24AA"),
]):
    y = 1.05 + i * 1.4
    c = RGBColor(int(col_hex[1:3],16), int(col_hex[3:5],16), int(col_hex[5:7],16))
    add_rect(sl11, 9.95, y, 3.2, 1.2, fill_color=RGBColor(0x0A, 0x1E, 0x40))
    add_rect(sl11, 9.95, y, 3.2, 0.07, fill_color=c)
    add_text(sl11, method, 10.05, y + 0.12, 3.0, 0.38,
             font_size=Pt(12), font_color=RGBColor(0xAA, 0xBB, 0xCC), bold=False)
    add_text(sl11, val, 10.05, y + 0.5, 3.0, 0.55,
             font_size=Pt(26), font_color=c, bold=True,
             align=PP_ALIGN.LEFT, font_name="Calibri Light")

add_rect(sl11, 9.95, 5.25, 3.2, 0.04, fill_color=GOLD)
add_text(sl11, f"Portfolio: ${final_val:,.0f}\n+{total_ret:.1f}% over 5 years",
         9.95, 5.35, 3.2, 0.75,
         font_size=Pt(13), font_color=WHITE, bold=False, align=PP_ALIGN.CENTER)
add_text(sl11, "95% confidence · 1-day horizon",
         9.95, 6.1, 3.2, 0.35,
         font_size=Pt(10), font_color=GOLD, align=PP_ALIGN.CENTER)

# =============================================================================
# SAVE
# =============================================================================

os.makedirs(OUT_FOLDER, exist_ok=True)
prs.save(OUT_FILE)
print(f"\nPresentation saved to:\n  {OUT_FILE}")
print("Open the file in PowerPoint or LibreOffice Impress.")
